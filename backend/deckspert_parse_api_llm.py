# DeckSpert.AI — FastAPI Parser API (LLM-Enhanced)
# ------------------------------------------------------------
# Endpoints
#   - POST /api/parse       : quick metrics + language + keywords (LLM-boosted when available)
#   - POST /api/sections    : structured IR sections extracted by LLM (with heuristics fallback)
#
# Behavior
#   - If OPENAI_API_KEY is present, uses LLM to: detect language, extract sections, refine metrics
#   - Otherwise falls back to lightweight heuristic parser
#
# Setup
#   pip install fastapi uvicorn pydantic langdetect python-pptx pypdf openai tiktoken
#   export OPENAI_API_KEY=sk-...
#   uvicorn deckspert_parse_api_llm:app --reload --port 8000

from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel, Field, validator
from typing import Dict, Any, List, Optional
import os, io, re, json

# Optional deps
try:
    from langdetect import detect
except Exception:
    detect = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None
try:
    from PyPDF2 import PdfReader
except Exception:
    PdfReader = None

# OpenAI LLM (optional)
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
USE_LLM = bool(OPENAI_API_KEY)
try:
    if USE_LLM:
        from openai import OpenAI
        oai = OpenAI(api_key=OPENAI_API_KEY)
except Exception:
    USE_LLM = False
    oai = None

app = FastAPI(title="DeckSpert.AI Parser API (LLM)", version="0.2.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ----------------------------- Models ----------------------------- #
class IRMetrics(BaseModel):
    clarity: float = Field(..., ge=0.0, le=1.0)
    logic: float = Field(..., ge=0.0, le=1.0)
    storytelling: float = Field(..., ge=0.0, le=1.0)
    data_strength: float = Field(..., ge=0.0, le=1.0)

class IRSections(BaseModel):
    problem: Optional[str] = ""
    solution: Optional[str] = ""
    market: Optional[str] = ""
    traction: Optional[str] = ""
    team: Optional[str] = ""
    business_model: Optional[str] = ""
    ask: Optional[str] = ""
    risks: Optional[str] = ""

class ParseResponse(BaseModel):
    language: str
    keywords: List[str]
    metrics: IRMetrics

class SectionsResponse(BaseModel):
    language: str
    sections: IRSections
    metrics: IRMetrics


# ----------------------------- Utils ----------------------------- #
def _read_txt_bytes(b: bytes) -> str:
    try:
        return b.decode("utf-8", errors="ignore")
    except Exception:
        return b.decode(errors="ignore")

def _read_pdf_bytes(b: bytes) -> str:
    if PdfReader is None:
        return ""
    try:
        reader = PdfReader(io.BytesIO(b))
        texts = []
        for page in reader.pages:
            try:
                texts.append(page.extract_text() or "")
            except Exception:
                continue
        return "\n".join(texts)
    except Exception:
        return ""

def _read_pptx_bytes(b: bytes) -> str:
    if Presentation is None:
        return ""
    try:
        prs = Presentation(io.BytesIO(b))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""


def quick_keywords(text: str, topk: int = 7) -> List[str]:
    if not text:
        return ["IR", "AI", "SaaS"]
    cleaned = re.sub(r"[^\w\s가-힣]", " ", text, flags=re.UNICODE).lower()
    tokens = [t for t in cleaned.split() if len(t) > 2]
    freq: Dict[str, int] = {}
    for t in tokens:
        freq[t] = freq.get(t, 1) + 1
    return [w for w, _ in sorted(freq.items(), key=lambda x: x[1], reverse=True)[:topk]]


def score_metrics_heuristic(text: str) -> IRMetrics:
    if not text:
        return IRMetrics(clarity=0.76, logic=0.72, storytelling=0.74, data_strength=0.70)
    # heuristic features
    length = len(text)
    bullets = len(re.findall(r"\n[-•·]", text))
    numbers = len(re.findall(r"\b\d+[\.,]?\d*\b", text))
    headings = len(re.findall(r"\n[A-Z가-힣].{0,60}\n", text))
    norm_len = min(1.0, max(0.0, length / 20000))
    norm_bullets = min(1.0, bullets / 40)
    norm_numbers = min(1.0, numbers / 80)
    norm_headings = min(1.0, headings / 30)
    clarity = 0.65 + 0.2 * norm_headings + 0.15 * norm_bullets
    logic = 0.62 + 0.25 * norm_headings + 0.15 * norm_len
    storytelling = 0.63 + 0.25 * norm_len + 0.12 * norm_bullets
    data_strength = 0.60 + 0.30 * norm_numbers + 0.10 * norm_len
    # clamp
    def c(v: float) -> float:
        return float(max(0.0, min(1.0, round(v, 4))))
    return IRMetrics(
        clarity=c(clarity), logic=c(logic), storytelling=c(storytelling), data_strength=c(data_strength)
    )


def detect_language(text: str, fallback: str = "en") -> str:
    if USE_LLM and text:
        try:
            prompt = "Detect the primary language (BCP-47 code like en, ko, ja, zh) of the following text. Answer with code only.\n\n" + text[:2000]
            r = oai.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                temperature=0
            )
            code = (r.choices[0].message.content or "").strip().lower()
            if re.match(r"^[a-z]{2}(-[a-z0-9]+)?$", code):
                return code
        except Exception:
            pass
    if detect is not None and text:
        try:
            return detect(text[:4000]) or fallback
        except Exception:
            return fallback
    return fallback


# ----------------------------- LLM Calls ----------------------------- #
LLM_SECTIONS_SYS = (
    "You are an expert VC analyst. Extract IR sections and score quality. "
    "Return STRICT JSON with keys: sections{problem,solution,market,traction,team,business_model,ask,risks}, "
    "metrics{clarity,logic,storytelling,data_strength} in 0..1."
)

LLM_SECTIONS_USER = (
    "Extract from the following content. Be concise. If info is missing, put an empty string. "
    "Text:\n\n{TEXT}"
)

def llm_extract_sections(text: str) -> Optional[SectionsResponse]:
    if not (USE_LLM and text):
        return None
    try:
        r = oai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": LLM_SECTIONS_SYS},
                {"role": "user", "content": LLM_SECTIONS_USER.replace("{TEXT}", text[:12000])},
            ],
            temperature=0.2,
            response_format={"type": "json_object"},
        )
        raw = r.choices[0].message.content or "{}"
        data = json.loads(raw)
        sections = IRSections(**data.get("sections", {}))
        m = data.get("metrics", {})
        metrics = IRMetrics(
            clarity=float(min(max(m.get("clarity", 0.76), 0.0), 1.0)),
            logic=float(min(max(m.get("logic", 0.72), 0.0), 1.0)),
            storytelling=float(min(max(m.get("storytelling", 0.74), 0.0), 1.0)),
            data_strength=float(min(max(m.get("data_strength", 0.70), 0.0), 1.0)),
        )
        lang = detect_language(text)
        return SectionsResponse(language=lang, sections=sections, metrics=metrics)
    except Exception:
        return None


# ----------------------------- Endpoints ----------------------------- #
@app.get("/health")
async def health() -> Dict[str, str]:
    return {"status": "ok", "llm": str(USE_LLM)}


@app.post("/api/parse")
async def parse_ir(file: UploadFile = File(...)) -> Any:
    try:
        raw = await file.read()
        name = (file.filename or "").lower()
        if name.endswith(".pdf"):
            text = _read_pdf_bytes(raw)
        elif name.endswith(".pptx") or name.endswith(".ppt"):
            text = _read_pptx_bytes(raw)
        elif name.endswith(".md") or name.endswith(".txt"):
            text = _read_txt_bytes(raw)
        else:
            text = _read_txt_bytes(raw)

        # Language
        lang = detect_language(text)

        # Keywords
        kws = quick_keywords(text)

        # Base metrics (heuristic)
        metrics = score_metrics_heuristic(text)

        # If LLM available, refine metrics slightly based on section completeness
        if USE_LLM and text:
            sec = llm_extract_sections(text)
            if sec:
                # adjust metrics by presence/length of key sections
                comp = sum(bool(getattr(sec.sections, k)) for k in [
                    "problem","solution","market","traction","team","business_model","ask" 
                ])
                bonus = min(0.06, comp * 0.008)  # up to +0.06
                metrics = IRMetrics(
                    clarity=min(1.0, metrics.clarity + bonus),
                    logic=min(1.0, metrics.logic + bonus),
                    storytelling=min(1.0, metrics.storytelling + bonus/2),
                    data_strength=min(1.0, metrics.data_strength + bonus/2),
                )

        payload = ParseResponse(language=lang, keywords=kws, metrics=metrics)
        return JSONResponse(payload.dict())

    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/sections")
async def parse_sections(file: UploadFile = File(...)) -> Any:
    try:
        raw = await file.read()
        name = (file.filename or "").lower()
        if name.endswith(".pdf"):
            text = _read_pdf_bytes(raw)
        elif name.endswith(".pptx") or name.endswith(".ppt"):
            text = _read_pptx_bytes(raw)
        elif name.endswith(".md") or name.endswith(".txt"):
            text = _read_txt_bytes(raw)
        else:
            text = _read_txt_bytes(raw)

        lang = detect_language(text)

        # Try LLM first
        sec = llm_extract_sections(text)
        if sec:
            return JSONResponse(sec.dict())

        # Fallback: very rough section slicing by keywords
        def slice_section(key: str, text: str) -> str:
            pattern = re.compile(rf"(^|\n)\s*{key}\s*[:\-]?\s*(.+?)(?=\n[A-Z가-힣][^\n]{0,40}\n|$)", re.IGNORECASE | re.DOTALL)
            m = pattern.search(text)
            return (m.group(2).strip() if m else "")[:1200]

        sections = IRSections(
            problem=slice_section("problem|문제", text),
            solution=slice_section("solution|해결", text),
            market=slice_section("market|시장", text),
            traction=slice_section("traction|지표", text),
            team=slice_section("team|팀", text),
            business_model=slice_section("business model|revenue|BM|수익", text),
            ask=slice_section("ask|funding|investment|라운드", text),
            risks=slice_section("risk|리스크", text),
        )
        metrics = score_metrics_heuristic(text)
        resp = SectionsResponse(language=lang, sections=sections, metrics=metrics)
        return JSONResponse(resp.dict())

    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

